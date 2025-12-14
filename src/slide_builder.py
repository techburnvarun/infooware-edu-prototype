# src/slide_builder.py
"""
Slide Builder Module
"""

import os
from pathlib import Path
from typing import List, Dict
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def build_pptx(slide_data: List[Dict], output_path: str, temp_frames_dir: str = None):
    """
    Build the PowerPoint presentation and optionally export slide images.
    """
    if not slide_data:
        raise ValueError("No slide data provided to build PPTX.")

    prs = Presentation()
    
    # Set slide size to standard 16:9 (default in recent PowerPoint)
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # Use a clean blank layout for full control
    blank_layout = prs.slide_layouts[6]  # Blank slide

    print(f"   Building {len(slide_data)} slides...")

    for idx, slide in enumerate(slide_data):
        sld = prs.slides.add_slide(blank_layout)

        # Background: white
        background = sld.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(255, 255, 255)

        title = slide["title"]
        bullets = slide.get("bullets", [])
        visual_path = slide.get("visual_path", "")

        # Dimensions and margins
        left_margin = Inches(0.8)
        top_margin = Inches(0.7)
        content_width = Inches(11.7)
        image_width = Inches(5.5)
        image_height = Inches(5.0)

        # Add title
        title_top = Inches(0.5)
        title_left = left_margin
        title_width = content_width
        title_height = Inches(1.2)

        title_box = sld.shapes.add_textbox(title_left, title_top, title_width, title_height)
        tf = title_box.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.1)
        tf.margin_top = Inches(0.1)

        p = tf.add_paragraph()
        p.text = title
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = RGBColor(33, 33, 33)
        p.alignment = PP_ALIGN.LEFT

        # Add bullets (left side)
        bullets_left = left_margin
        bullets_top = Inches(2.0)
        bullets_width = Inches(6.0)
        bullets_height = Inches(4.5)

        bullets_box = sld.shapes.add_textbox(bullets_left, bullets_top, bullets_width, bullets_height)
        bf = bullets_box.text_frame
        bf.word_wrap = True
        bf.margin_left = Inches(0.2)
        bf.margin_top = Inches(0.1)

        for bullet in bullets:
            p = bf.add_paragraph()
            # Remove leading • if already present
            bullet_text = bullet.strip().lstrip("• ").strip()
            p.text = "• " + bullet_text
            p.font.size = Pt(24)
            p.font.color.rgb = RGBColor(50, 50, 50)
            p.space_after = Pt(12)
            p.level = 0

        # Add image (right side)
        if visual_path and os.path.exists(visual_path):
            img_left = Inches(7.2)
            img_top = Inches(1.5)

            try:
                sld.shapes.add_picture(
                    visual_path,
                    img_left,
                    img_top,
                    width=image_width,
                    height=image_height
                )
            except Exception as e:
                print(f"   Warning: Could not add image for slide {idx+1}: {e}")
        else:
            # Optional: add placeholder text if no image
            pass

        # Progress
        print(f"      Slide {idx+1}/{len(slide_data)}: '{title[:60]}...'")

    # Save the PPTX
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    prs.save(output_path)

    # Export slides as PNG images if temp_frames_dir is provided
    if temp_frames_dir:
        temp_frames_dir = os.path.abspath(temp_frames_dir)
        os.makedirs(temp_frames_dir, exist_ok=True)

        print(f"   Exporting slides to {temp_frames_dir} using PowerPoint...")

        try:
            import comtypes.client

            powerpoint = comtypes.client.CreateObject("PowerPoint.Application")
            powerpoint.Visible = 1  # Set to 0 once stable

            presentation = powerpoint.Presentations.Open(
            os.path.abspath(output_path),
            WithWindow=False
            )

            slide_count = presentation.Slides.Count

            for i in range(1, slide_count + 1):  # PowerPoint is 1-based
                frame_path = os.path.abspath(
                    os.path.join(temp_frames_dir, f"slide_{i:03d}.png")
                )

                presentation.Slides(i).Export(frame_path, "PNG", 1920, 1080)
                print(f"      Exported frame: {os.path.basename(frame_path)}")

            presentation.Close()
            powerpoint.Quit()

            print(f"   Successfully exported {slide_count} PNG frames")

        except Exception as e:
            print(f"   ERROR: Failed to export slides via PowerPoint COM: {e}")
            print("   Make sure Microsoft PowerPoint is installed and accessible.")
            raise
